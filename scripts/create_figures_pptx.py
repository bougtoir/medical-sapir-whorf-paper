#!/usr/bin/env python3
"""
Generate editable PPTX with figures and tables for:
"The Nosological Imperative" manuscript

Output: medical_sapir_whorf/output/figures_tables_en.pptx
        medical_sapir_whorf/output/figures_tables_ja.pptx
"""

import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def draw_ntl_diagram(slide, lang="en"):
    """Draw the Nosological Therapeutic Loop diagram on the slide."""

    if lang == "en":
        labels = [
            ("N", "Nosological\nAct", "Disease category\ncreated/refined"),
            ("V", "Patient\nValidation", "Name reduces\ndiagnostic uncertainty"),
            ("A", "Clinical\nAttention", "New cognitive\ncategories for clinicians"),
            ("R", "Research\nInvestment", "Funding & trials\norganized around category"),
            ("T", "Treatment\nDevelopment", "New therapeutics\nfor named condition"),
            ("N\u2032", "Nosological\nRefinement", "Category revised\nbased on evidence"),
        ]
        boundary_label = "BOUNDARY\nCONDITION"
        boundary_desc = "Without empirical\ngrounding \u2192\nDisease Mongering"
        caption = (
            "Figure 1. The Nosological Therapeutic Loop (NTL). "
            "Disease naming triggers a self-reinforcing cycle through six stages."
        )
    else:
        labels = [
            ("N", "\u75be\u60a3\u547d\u540d\n\u884c\u70ba", "\u75be\u60a3\u30ab\u30c6\u30b4\u30ea\u30fc\u306e\n\u5275\u8a2d\u30fb\u7cbe\u7dfb\u5316"),
            ("V", "\u60a3\u8005\u306e\n\u627f\u8a8d", "\u8a3a\u65ad\u7684\u4e0d\u78ba\u5b9f\u6027\u306e\n\u8efd\u6e1b"),
            ("A", "\u81e8\u5e8a\u7684\n\u6ce8\u76ee", "\u65b0\u305f\u306a\u8a8d\u77e5\n\u30ab\u30c6\u30b4\u30ea\u30fc\u306e\u7372\u5f97"),
            ("R", "\u7814\u7a76\n\u6295\u8cc7", "\u8cc7\u91d1\u30fb\u8a66\u9a13\u306e\n\u30ab\u30c6\u30b4\u30ea\u30fc\u5468\u8fba\u3078\u306e\u7d44\u7e54\u5316"),
            ("T", "\u6cbb\u7642\n\u958b\u767a", "\u547d\u540d\u3055\u308c\u305f\u72b6\u614b\u3078\u306e\n\u65b0\u6cbb\u7642\u6cd5"),
            ("N\u2032", "\u75be\u60a3\u547d\u540d\u306e\n\u7cbe\u7dfb\u5316", "\u30a8\u30d3\u30c7\u30f3\u30b9\u306b\u57fa\u3065\u304f\n\u30ab\u30c6\u30b4\u30ea\u30fc\u6539\u8a02"),
        ]
        boundary_label = "\u5883\u754c\u6761\u4ef6"
        boundary_desc = "\u5b9f\u8a3c\u7684\u6839\u62e0\u306a\u304d\u5834\u5408\n\u2192 \u75be\u60a3\u5587\u4f1d"
        caption = (
            "\u56f31. \u75be\u60a3\u547d\u540d\u6cbb\u7642\u30eb\u30fc\u30d7\uff08NTL\uff09\u3002"
            "\u75be\u60a3\u547d\u540d\u304c6\u3064\u306e\u6bb5\u968e\u3092\u7d4c\u308b\u81ea\u5df1\u5f37\u5316\u30b5\u30a4\u30af\u30eb\u3092\u8a98\u767a\u3059\u308b\u3002"
        )

    # Circle layout parameters
    cx, cy = Inches(6.0), Inches(3.4)
    radius = Inches(2.2)
    node_w, node_h = Inches(1.8), Inches(1.3)
    desc_w, desc_h = Inches(1.9), Inches(0.8)

    colors = [
        RGBColor(41, 98, 255),    # N - blue
        RGBColor(0, 150, 136),    # V - teal
        RGBColor(76, 175, 80),    # A - green
        RGBColor(255, 152, 0),    # R - orange
        RGBColor(233, 30, 99),    # T - pink
        RGBColor(103, 58, 183),   # N' - purple
    ]

    n = len(labels)
    positions = []
    for i in range(n):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + radius * math.cos(angle)
        y = cy + radius * math.sin(angle)
        positions.append((x, y))

    # Draw arrows between nodes
    for i in range(n):
        x1, y1 = positions[i]
        x2, y2 = positions[(i + 1) % n]
        # Shorten arrows
        dx, dy = x2 - x1, y2 - y1
        dist = math.sqrt(dx**2 + dy**2)
        shorten = Inches(0.9)
        if dist > 0:
            ratio = shorten / dist
            ax1 = x1 + dx * ratio
            ay1 = y1 + dy * ratio
            ax2 = x2 - dx * ratio
            ay2 = y2 - dy * ratio

            connector = slide.shapes.add_connector(
                1,  # straight connector
                int(ax1), int(ay1),
                int(ax2), int(ay2)
            )
            connector.line.color.rgb = RGBColor(150, 150, 150)
            connector.line.width = Pt(2)

    # Draw nodes
    for i, ((abbrev, label, desc), (x, y), color) in enumerate(
        zip(labels, positions, colors)
    ):
        # Main node circle
        left = int(x - node_w / 2)
        top = int(y - node_h / 2)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, int(node_w), int(node_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Abbreviation
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = abbrev
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Label
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = label
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(255, 255, 255)

        # Description box (outside)
        angle = -math.pi / 2 + 2 * math.pi * i / n
        desc_r = radius + Inches(1.6)
        dx = cx + desc_r * math.cos(angle)
        dy = cy + desc_r * math.sin(angle)
        desc_left = int(dx - desc_w / 2)
        desc_top = int(dy - desc_h / 2)
        desc_box = slide.shapes.add_textbox(desc_left, desc_top, int(desc_w), int(desc_h))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(9)
        dp.font.color.rgb = RGBColor(80, 80, 80)
        dp.alignment = PP_ALIGN.CENTER

    # Boundary condition box (center)
    bc_w, bc_h = Inches(2.0), Inches(1.4)
    bc_left = int(cx - bc_w / 2)
    bc_top = int(cy - bc_h / 2)
    bc_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        bc_left, bc_top, int(bc_w), int(bc_h)
    )
    bc_shape.fill.solid()
    bc_shape.fill.fore_color.rgb = RGBColor(244, 67, 54)
    bc_shape.line.fill.background()

    btf = bc_shape.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = boundary_label
    br.font.size = Pt(12)
    br.font.bold = True
    br.font.color.rgb = RGBColor(255, 255, 255)

    bp2 = btf.add_paragraph()
    bp2.alignment = PP_ALIGN.CENTER
    br2 = bp2.add_run()
    br2.text = boundary_desc
    br2.font.size = Pt(9)
    br2.font.color.rgb = RGBColor(255, 220, 220)

    # Caption at bottom
    cap_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.6)
    )
    ctf = cap_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(10)
    cp.font.italic = True
    cp.alignment = PP_ALIGN.LEFT


def add_table_slide(prs, title, caption, headers, rows, lang="en"):
    """Add a table slide to the presentation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Table
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.2)
    tbl_width = Inches(12.333)
    tbl_height = Inches(4.5)

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height).table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(41, 98, 255)

    # Data rows
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 250)

    # Caption at bottom
    cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8))
    ctf = cap_box.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(10)
    cp.font.italic = True
    cp.alignment = PP_ALIGN.LEFT


def create_pptx_en():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    add_title_slide(
        prs,
        "The Nosological Imperative",
        "Figures and Tables — English Version"
    )

    # Figure 1: NTL diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Figure 1. The Nosological Therapeutic Loop (NTL)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    draw_ntl_diagram(slide, lang="en")

    # Table 1
    add_table_slide(
        prs,
        title="Table 1. ICD-11 Chronic Pain Classification (MG30)",
        caption=(
            "Table 1. ICD-11 Chronic Pain Classification (MG30) with Medical "
            "Sapir-Whorf implications. Each category represents a nosological "
            "innovation over ICD-10 with specific implications for clinical "
            "cognition and patient experience."
        ),
        headers=["ICD-11 Code", "Category", "Key Innovation over ICD-10", "Sapir-Whorf Implication"],
        rows=[
            ["MG30.0", "Chronic primary pain", "Pain as disease entity, not symptom",
             "Validates pain without identifiable cause"],
            ["MG30.1", "Chronic cancer-related pain", "Distinguishes cancer pain from cancer",
             "Enables pain-specific treatment alongside oncology"],
            ["MG30.2", "Chronic postsurgical/post-traumatic pain", "Previously invisible category",
             "Makes iatrogenic chronification visible"],
            ["MG30.3", "Chronic secondary MSK pain", "Differentiates chronic from acute MSK",
             "Shifts framing to pain chronification"],
            ["MG30.4", "Chronic secondary visceral pain", "Previously organ-specific only",
             "Recognizes pain as independent clinical target"],
            ["MG30.5", "Chronic neuropathic pain", "Unified category across etiologies",
             "Creates coherent research framework"],
            ["MG30.6", "Chronic secondary headache/orofacial", "Chronic-specific classification",
             "Distinguishes chronic management from acute"],
        ]
    )

    # Table 2
    add_table_slide(
        prs,
        title="Table 2. Therapeutic Naming vs. Disease Mongering",
        caption=(
            "Table 2. Distinguishing therapeutic naming from disease mongering "
            "along three dimensions: origin, evidence base, and outcome trajectory."
        ),
        headers=["Dimension", "Therapeutic Naming", "Disease Mongering"],
        rows=[
            ["Origin",
             "Driven by clinical need and patient suffering; expert consensus",
             "Driven by commercial interest; pharmaceutical marketing"],
            ["Evidence base",
             "Grounded in empirical observation; validated through field testing",
             "Based on medicalization of normal variation; limited validation"],
            ["Outcome trajectory",
             "Improved patient outcomes, better resource allocation (NTL positive cycle)",
             "Overdiagnosis, unnecessary treatment, iatrogenic harm"],
        ]
    )

    output_path = os.path.join(OUTPUT_DIR, "figures_tables_en.pptx")
    prs.save(output_path)
    print(f"PPTX (EN) saved to: {output_path}")
    return output_path


def create_pptx_ja():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Title slide
    add_title_slide(
        prs,
        "\u75be\u60a3\u547d\u540d\u306e\u6cbb\u7642\u7684\u610f\u7fa9",
        "\u56f3\u8868 \u2014 \u65e5\u672c\u8a9e\u7248"
    )

    # Figure 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "\u56f31. \u75be\u60a3\u547d\u540d\u6cbb\u7642\u30eb\u30fc\u30d7\uff08NTL\uff09"
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    draw_ntl_diagram(slide, lang="ja")

    # Table 1
    add_table_slide(
        prs,
        title="\u88681. ICD-11\u6162\u6027\u75bc\u75db\u5206\u985e\uff08MG30\uff09",
        caption=(
            "\u88681. ICD-11\u6162\u6027\u75bc\u75db\u5206\u985e\uff08MG30\uff09\u3068\u533b\u7642\u7248\u30b5\u30d4\u30a2\uff1d\u30a6\u30a9\u30fc\u30d5\u4eee\u8aac\u306e\u542b\u610f\u3002"
            "\u5404\u30ab\u30c6\u30b4\u30ea\u30fc\u306fICD-10\u304b\u3089\u306e\u75be\u60a3\u547d\u540d\u7684\u9769\u65b0\u3068\u305d\u306e\u81e8\u5e8a\u8a8d\u77e5\u30fb\u60a3\u8005\u7d4c\u9a13\u3078\u306e\u542b\u610f\u3092\u793a\u3059\u3002"
        ),
        headers=["ICD-11\u30b3\u30fc\u30c9", "\u30ab\u30c6\u30b4\u30ea\u30fc", "ICD-10\u304b\u3089\u306e\u9769\u65b0", "\u30b5\u30d4\u30a2\uff1d\u30a6\u30a9\u30fc\u30d5\u7684\u542b\u610f"],
        rows=[
            ["MG30.0", "\u6162\u6027\u4e00\u6b21\u6027\u75bc\u75db", "\u75bc\u75db\u3092\u75c7\u72b6\u3067\u306f\u306a\u304f\u75be\u60a3\u3068\u3057\u3066\u8a8d\u5b9a", "\u539f\u56e0\u4e0d\u660e\u306e\u75bc\u75db\u3092\u627f\u8a8d"],
            ["MG30.1", "\u6162\u6027\u304c\u3093\u95a2\u9023\u75bc\u75db", "\u304c\u3093\u306e\u75bc\u75db\u3092\u304c\u3093\u305d\u306e\u3082\u306e\u304b\u3089\u533a\u5225", "\u75bc\u75db\u7279\u7570\u7684\u6cbb\u7642\u3092\u53ef\u80fd\u306b"],
            ["MG30.2", "\u6162\u6027\u8853\u5f8c\u30fb\u5916\u50b7\u5f8c\u75bc\u75db", "\u4ee5\u524d\u306f\u4e0d\u53ef\u8996\u306e\u30ab\u30c6\u30b4\u30ea\u30fc", "\u533b\u539f\u6027\u6162\u6027\u5316\u3092\u53ef\u8996\u5316"],
            ["MG30.3", "\u6162\u6027\u4e8c\u6b21\u6027\u7b4b\u9aa8\u683c\u7cfb\u75bc\u75db", "\u6162\u6027\u3068\u6025\u6027\u3092\u533a\u5225", "\u75bc\u75db\u6162\u6027\u5316\u3078\u306e\u67a0\u7d44\u307f\u8ee2\u63db"],
            ["MG30.4", "\u6162\u6027\u4e8c\u6b21\u6027\u5185\u81d3\u75db", "\u4ee5\u524d\u306f\u81d3\u5668\u7279\u7570\u7684\u8a3a\u65ad\u306e\u307f", "\u75bc\u75db\u3092\u72ec\u7acb\u3057\u305f\u81e8\u5e8a\u30bf\u30fc\u30b2\u30c3\u30c8\u3068\u3057\u3066\u8a8d\u8b58"],
            ["MG30.5", "\u6162\u6027\u795e\u7d4c\u969c\u5bb3\u6027\u75bc\u75db", "\u75c5\u56e0\u6a2a\u65ad\u7684\u306a\u7d71\u4e00\u30ab\u30c6\u30b4\u30ea\u30fc", "\u4e00\u8cab\u3057\u305f\u7814\u7a76\u67a0\u7d44\u307f\u3092\u5275\u51fa"],
            ["MG30.6", "\u6162\u6027\u4e8c\u6b21\u6027\u982d\u75db\u30fb\u53e3\u8154\u9854\u9762\u75db", "\u6162\u6027\u7279\u7570\u7684\u5206\u985e", "\u6162\u6027\u7ba1\u7406\u3092\u6025\u6027\u304b\u3089\u533a\u5225"],
        ]
    )

    # Table 2
    add_table_slide(
        prs,
        title="\u88682. \u6cbb\u7642\u7684\u547d\u540d vs. \u75be\u60a3\u5587\u4f1d",
        caption=(
            "\u88682. \u6cbb\u7642\u7684\u547d\u540d\u3068\u75be\u60a3\u5587\u4f1d\u306e\u533a\u5225\u3002"
            "\u8d77\u6e90\u3001\u30a8\u30d3\u30c7\u30f3\u30b9\u57fa\u76e4\u3001\u30a2\u30a6\u30c8\u30ab\u30e0\u8ecc\u9053\u306e\u4e09\u6b21\u5143\u306b\u6cbf\u3063\u305f\u6bd4\u8f03\u3002"
        ),
        headers=["\u6b21\u5143", "\u6cbb\u7642\u7684\u547d\u540d", "\u75be\u60a3\u5587\u4f1d"],
        rows=[
            ["\u8d77\u6e90", "\u81e8\u5e8a\u7684\u5fc5\u8981\u6027\u3068\u60a3\u8005\u306e\u82e6\u75db\u304c\u99c6\u52d5\uff1b\u5c02\u9580\u5bb6\u30b3\u30f3\u30bb\u30f3\u30b5\u30b9",
             "\u5546\u696d\u7684\u5229\u76ca\u304c\u99c6\u52d5\uff1b\u88fd\u85ac\u30de\u30fc\u30b1\u30c6\u30a3\u30f3\u30b0"],
            ["\u30a8\u30d3\u30c7\u30f3\u30b9\u57fa\u76e4", "\u7d4c\u9a13\u7684\u89b3\u5bdf\u306b\u57fa\u790e\u3065\u3051\uff1b\u30d5\u30a3\u30fc\u30eb\u30c9\u30c6\u30b9\u30c8\u3067\u691c\u8a3c",
             "\u6b63\u5e38\u5909\u7570\u306e\u533b\u7642\u5316\uff1b\u72ec\u7acb\u3057\u305f\u691c\u8a3c\u304c\u9650\u5b9a\u7684"],
            ["\u30a2\u30a6\u30c8\u30ab\u30e0\u8ecc\u9053", "\u60a3\u8005\u30a2\u30a6\u30c8\u30ab\u30e0\u306e\u6539\u5584\u3001\u8cc7\u6e90\u914d\u5206\u306e\u6539\u5584\uff08NTL\u6b63\u306e\u30b5\u30a4\u30af\u30eb\uff09",
             "\u904e\u5270\u8a3a\u65ad\u3001\u4e0d\u5fc5\u8981\u306a\u6cbb\u7642\u3001\u533b\u539f\u6027\u50b7\u5bb3"],
        ]
    )

    output_path = os.path.join(OUTPUT_DIR, "figures_tables_ja.pptx")
    prs.save(output_path)
    print(f"PPTX (JA) saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    create_pptx_en()
    create_pptx_ja()
